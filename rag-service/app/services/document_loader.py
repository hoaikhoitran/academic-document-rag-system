"""
app/services/document_loader.py
===============================

STEP 1 of the RAG pipeline: read a file from disk and turn it into a
list of plain-text "pages" (or "sections" when the file format has no
real concept of a page, e.g. DOCX / TXT).

The output of this module is ALWAYS the same shape, regardless of the
input file type:

    [
        {"text": "...", "pageNumber": 1, "source": "file.pdf"},
        {"text": "...", "pageNumber": 2, "source": "file.pdf"},
        ...
    ]

Why this matters for RAG
------------------------
Vector databases store small chunks of text together with metadata.
By keeping page numbers here (not invented later), we can show the
student WHICH page of the PDF a sentence came from when we cite
sources in the final answer.

Supported file types
--------------------
* PDF  -> PyMuPDF (fitz)
* DOCX -> python-docx
* PPTX -> python-pptx       (each slide is treated as one "page")
* TXT  -> built-in `open()`
"""

from __future__ import annotations

import os
from pathlib import Path
from typing import Any

from app.utils.file_utils import get_file_extension, is_supported_file
from app.utils.text_utils import normalize_whitespace


# ---------------------------------------------------------------------------
# Internal loaders. Each one returns the same list-of-dicts shape so the
# caller never needs to care about the original file format.
# ---------------------------------------------------------------------------


def _load_pdf(file_path: str) -> list[dict[str, Any]]:
    """
    Read a PDF using PyMuPDF (imported as `fitz`).

    PyMuPDF is fast and preserves page boundaries naturally — we can
    iterate `doc` and call `page.get_text("text")` to get the text of
    each page in reading order.
    """
    import fitz  # PyMuPDF

    source = os.path.basename(file_path)
    pages: list[dict[str, Any]] = []

    with fitz.open(file_path) as doc:
        for index, page in enumerate(doc):
            raw = page.get_text("text") or ""
            cleaned = normalize_whitespace(raw)
            if not cleaned:
                # Skip empty pages (e.g. cover pages, images-only pages).
                continue
            pages.append(
                {
                    "text": cleaned,
                    "pageNumber": index + 1,  # 1-based, like a human counts pages
                    "source": source,
                }
            )

    return pages


def _load_docx(file_path: str) -> list[dict[str, Any]]:
    """
    Read a DOCX (Word) file.

    DOCX has no real "pages" concept (pagination is decided by the word
    processor at render time), so we return ONE section containing the
    concatenation of every paragraph. pageNumber is None.
    """
    from docx import Document  # python-docx

    source = os.path.basename(file_path)
    document = Document(file_path)

    paragraphs = [p.text for p in document.paragraphs if p.text and p.text.strip()]
    cleaned = normalize_whitespace("\n".join(paragraphs))

    if not cleaned:
        return []

    return [
        {
            "text": cleaned,
            "pageNumber": None,
            "source": source,
        }
    ]


def _load_pptx(file_path: str) -> list[dict[str, Any]]:
    """
    Read a PowerPoint file.

    Each SLIDE is treated as one "page" so we can show the student which
    slide an answer came from. We concatenate the text from every shape
    on the slide (titles, bullets, text boxes).
    """
    from pptx import Presentation  # python-pptx

    source = os.path.basename(file_path)
    presentation = Presentation(file_path)
    pages: list[dict[str, Any]] = []

    for index, slide in enumerate(presentation.slides):
        pieces: list[str] = []
        for shape in slide.shapes:
            # `has_text_frame` is True for titles, content placeholders,
            # and free text boxes — all the places students put info.
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text:
                            pieces.append(run.text)

        cleaned = normalize_whitespace(" ".join(pieces))
        if not cleaned:
            continue

        pages.append(
            {
                "text": cleaned,
                "pageNumber": index + 1,  # slide number, 1-based
                "source": source,
            }
        )

    return pages


def _load_txt(file_path: str) -> list[dict[str, Any]]:
    """
    Read a plain-text file. UTF-8 first, fall back to latin-1 if needed
    (some Windows-exported notes use cp1252 / latin-1).
    """
    source = os.path.basename(file_path)

    try:
        with open(file_path, "r", encoding="utf-8") as fh:
            raw = fh.read()
    except UnicodeDecodeError:
        with open(file_path, "r", encoding="latin-1") as fh:
            raw = fh.read()

    cleaned = normalize_whitespace(raw)
    if not cleaned:
        return []

    return [
        {
            "text": cleaned,
            "pageNumber": None,
            "source": source,
        }
    ]


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------


def load_document(file_path: str) -> list[dict[str, Any]]:
    """
    Read a document from disk and return a normalized list of pages.

    Raises
    ------
    FileNotFoundError
        If the file does not exist on disk.
    ValueError
        If the extension is not one of the supported formats.
    RuntimeError
        If the file exists but cannot be read (corrupted, password-protected, ...).
    """
    if not os.path.isfile(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    if not is_supported_file(file_path):
        raise ValueError(
            "Unsupported file type. Supported extensions: .pdf, .docx, .pptx, .txt"
        )

    extension = get_file_extension(file_path)

    try:
        if extension == ".pdf":
            return _load_pdf(file_path)
        if extension == ".docx":
            return _load_docx(file_path)
        if extension == ".pptx":
            return _load_pptx(file_path)
        if extension == ".txt":
            return _load_txt(file_path)
    except Exception as exc:  # noqa: BLE001 — re-raise as a clear RuntimeError below
        raise RuntimeError(f"Failed to read {Path(file_path).name}: {exc}") from exc

    # Should be unreachable because of the `is_supported_file` guard above.
    raise ValueError(f"Unhandled file extension: {extension}")
